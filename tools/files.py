"""Process a document and return its content.

Largely inspired from https://github.dev/aymeric-roucher/GAIA/blob/main/scripts/tools/mdconvert.py
"""

import re
import html
from typing import Union
import logging
from langchain_core.tools import tool
from bs4 import BeautifulSoup
import markdownify
import pdfplumber
import mammoth
import os
import pptx
import xml.etree.ElementTree as ET
import zipfile
import tempfile
import requests
import mimetypes
from urllib.parse import urlparse

logger = logging.getLogger(__name__)


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title = title
        self.text_content = text_content

    def __str__(self) -> str:
        if self.title:
            return f"{self.title}\n\n{self.text_content}"
        return self.text_content


class DocumentConverter:
    extensions: list[str] = []

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()

    def validate_extension(self, local_path) -> bool:
        extension = local_path.split(".")[-1]
        return extension.lower() in self.extensions


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    extensions: list[str] = ["html", "htm"]

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        if not self.validate_extension(local_path):
            return None

        result = None
        with open(local_path) as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = markdownify.MarkdownConverter().convert_soup(body_elm)
        else:
            webpage_text = markdownify.MarkdownConverter().convert_soup(soup)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    extensions: list[str] = [
        "txt",
        "py",
        "js",
        "css",
        "json",
        "md",
        "go",
        "rb",
        "rs",
        "rsx",
    ]

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        if not self.validate_extension(local_path):
            return None

        text_content = ""
        with open(local_path) as fh:
            text_content = fh.read()

        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class PdfConverter(DocumentConverter):
    extensions: list[str] = ["pdf"]

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        if not self.validate_extension(local_path):
            return None

        all_text = ""

        with pdfplumber.open(local_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    all_text += f"\n\n--- Page {i + 1} ---\n\n{text}"

        return DocumentConverterResult(
            title=None,
            text_content=markdownify.markdownify(all_text).strip(),
        )


class DocxConverter(HtmlConverter):
    extensions: list[str] = ["docx"]

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        if not self.validate_extension(local_path):
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class PptxConverter(HtmlConverter):
    extensions: list[str] = ["pptx"]

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        if not self.validate_extension(local_path):
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except:  # noqa
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    # try:
                    #    filename = shape.image.filename
                    # except:
                    #    pass

                    md_content += (
                        "\n!["
                        + (alt_text if alt_text else shape.name)
                        + "]("
                        + filename
                        + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + " "
                    else:
                        md_content += shape.text + " "

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False


class XmlConverter(DocumentConverter):
    extensions: list[str] = ["xml"]

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        if not self.validate_extension(local_path):
            return None

        xml_string = ""
        with open(local_path) as fh:
            xml_string = fh.read()

        def extract_table_from_html_like(xml_root):
            table = xml_root.find(".//table")
            if table is None:
                raise ValueError("No table found in the XML")

            headers = [th.text for th in table.find("thead").findall("th")]
            rows = [
                [td.text for td in tr.findall("td")]
                for tr in table.find("tbody").findall("tr")
            ]

            # Create markdown table
            markdown = "| " + " | ".join(headers) + " |\n"
            markdown += "| " + " | ".join(["---"] * len(headers)) + " |\n"
            for row in rows:
                markdown += "| " + " | ".join(row) + " |\n"

        def extract_table_from_wordml(xml_root, namespaces):
            # Parse the XML content
            root = xml_root
            namespace = {"w": "http://schemas.microsoft.com/office/word/2003/wordml"}

            # Extract text content
            body = root.find("w:body", namespace)
            paragraphs = body.findall(".//w:p", namespace)
            text_content = []
            for para in paragraphs:
                texts = para.findall(".//w:t", namespace)
                for text in texts:
                    text_content.append(text.text)

            return "\n".join(text_content)

        # Parse the XML string
        root = ET.fromstring(xml_string)
        namespaces = {"w": "http://schemas.microsoft.com/office/word/2003/wordml"}

        if root.tag.endswith("wordDocument"):
            markdown = extract_table_from_wordml(root, namespaces)
        else:
            markdown = extract_table_from_html_like(root)

        return DocumentConverterResult(
            title=None,
            text_content=markdown.strip(),
        )


class DocumentConverterFactory:
    _converters: dict[str, DocumentConverter] = {}

    def get_converter(cls, extension: str) -> DocumentConverter:
        return cls._converters.get(extension.lower())

    def register_converter(cls, converter: DocumentConverter):
        for extension in converter.extensions:
            cls._converters[extension.lower()] = converter


converter_factory = DocumentConverterFactory()
converter_factory.register_converter(HtmlConverter())
converter_factory.register_converter(PdfConverter())
converter_factory.register_converter(PlainTextConverter())
converter_factory.register_converter(DocxConverter())
converter_factory.register_converter(PptxConverter())
converter_factory.register_converter(XmlConverter())


def save_resource(url: str) -> str:
    """Save a resource from a URL to a temporary file.

    Args:
        url (str): The URL of the resource.

    Returns:
        str: The path to the temporary file.
    """
    print("Downloading resource from URL:", url)
    response = requests.get(url, stream=True)
    response.raise_for_status()

    # Determine filename from URL path
    parsed_url = urlparse(url)
    base_name = os.path.basename(parsed_url.path)

    # Guess extension from Content-Type header
    content_type = response.headers.get("Content-Type")
    extension = (
        mimetypes.guess_extension(content_type.split(";")[0]) if content_type else ""
    )

    if not extension and "." in base_name:
        # Fall back to using URL extension
        extension = os.path.splitext(base_name)[1]

    print("Extension guessed from response:", extension)

    with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as tmp_file:
        # Write content to temporary file
        for chunk in response.iter_content(chunk_size=8192):
            tmp_file.write(chunk)
        tmp_path = tmp_file.name
        print("File saved to: ", tmp_path)
        return tmp_path


@tool
def load_file_or_url(file_path_or_url: str) -> str:
    """Load a file or a URL and return its contents.

    Use it for PDF, DOCX, HTML, PPTX, XML and any text resource.

    Args:
        file_path_or_url (str): The path to the file or URL.

    Returns:
        str: The contents of the file.
    """
    content = ""
    file_path = file_path_or_url
    # if URL, download first as a temporary file
    if file_path_or_url.startswith("http"):
        file_path = save_resource(file_path_or_url)
        content += f"URL: {file_path_or_url}\n"
        content += f"Downloaded to: {file_path}\n"

    extension = file_path.split(".")[-1].lower()
    converter = converter_factory.get_converter(extension)
    if converter:
        result = converter.convert(file_path)
        if result:
            content += str(result)
            # limit content to 5000 characters
            if len(content) > 5000:
                return content[:5000] + "...TRUNCATED"
            return content

    print(f"ERROR: Unable to load file or URL: Unknown format. Extension: {extension}")
    return "Unable to load file or URL: Unknown format"


@tool
def unzip(file_path: str) -> list[str]:
    """Unzip a file and return the list of files in the zip.
    Always use this to process zip files.

    Args:
        file_path (str): The path to the zip file.

    Returns:
        list[str]: The list of files in the zip.
    """
    with zipfile.ZipFile(file_path) as zip_file:
        # extract all files from the zip
        zip_file.extractall("data/zip")
        filepaths = [os.path.join("data/zip", f) for f in zip_file.namelist()]
        print("Extracted files:", filepaths)
        return filepaths
